from typing import List, Dict, Any
from pptx import Presentation
from pathlib import Path
from langchain.text_splitter import RecursiveCharacterTextSplitter
from config.config import CHUNK_SIZE, CHUNK_OVERLAP

class PPTProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            length_function=len,
            is_separator_regex=False,
        )

    def extract_text(self, file_path: str | Path) -> str:
        """Extract text from PPT file."""
        text = ""
        prs = Presentation(file_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text

    def process_document(self, file_path: str | Path) -> List[Dict[str, Any]]:
        """Process PPT document and return chunks with metadata."""
        text = self.extract_text(file_path)
        chunks = self.text_splitter.split_text(text)
        
        # Add metadata to chunks
        processed_chunks = []
        for i, chunk in enumerate(chunks):
            processed_chunks.append({
                'text': chunk,
                'source': Path(file_path).name,
                'chunk_index': i
            })
        return processed_chunks 